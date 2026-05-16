#!/usr/bin/env python3
"""
Fitness Factory Hattingen — Relaunch-Präsentation v5
25 Slides im Karteikarten-Stil (PowerPoint)

Änderungen gegenüber v4:
 - 5 neue Vater-Slides (2.1) integriert:
   · NEU Slide 5: Navigation-Aufgabe ("3 Klicks bis zum Ziel — nicht 13 Links")
   · NEU Slide 7: Kursplan-Aufgabe ("Der Kursplan ist eines der wichtigsten Werkzeuge")
   · NEU Slide 11: Logo-Aufgabe ("Ein Logo muss in 1 Sekunde sagen …")
   · NEU Slide 12: Logo-Tiefenanalyse (4 Probleme: 3D-Metallic · ITNESS · Zielgruppe · Skalierbarkeit)
   · NEU Slide 13: Logo Option A vs B (Blitz bleibt / Neues Symbol)
 - v4-Slide "Farbwelt" gestrichen (Logo-Empfehlung wandert in die neuen Logo-Slides)
 - v4-Slide "Abschluss" mit "Nächste Schritte" fusioniert (Bereit?-Box trägt Claim)
 - Design-Akzente (dezent, Vater-nah):
   · rote 6pt vertikale Linie links auf dunklen Slides
   · rote Sektions-Labels oben (Caps)
   · Footer "X / 25"
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Paths ──────────────────────────────────────────────
BASE = os.path.dirname(os.path.abspath(__file__))
SCREENSHOTS_NEW = os.path.join(BASE, "figma-screenshots", "new")
SCREENSHOTS_OLD = os.path.join(BASE, "figma-screenshots", "old")
SCREENSHOTS_V4 = os.path.join(BASE, "figma-screenshots", "v4")
FF_SLIDES = os.path.join(BASE, "FitnessFactory")
SCHILDER = "/Users/lukekozik/Documents/Programme/n8n/Jobs/FitnessFactory/Schilder"
PROJECT_ROOT = os.path.abspath(os.path.join(BASE, ".."))
LOGO_PATH = os.path.join(PROJECT_ROOT, "public", "images", "logo-ff-jpg-cropped.jpg")
OUTPUT = os.path.join(BASE, "FitnessFactory_Relaunch_v5.pptx")
TOTAL_SLIDES = 25

# ── Brand Colors ───────────────────────────────────────
BLACK = RGBColor(0x11, 0x11, 0x11)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
DARKER = RGBColor(0x14, 0x14, 0x14)
RED = RGBColor(0xE7, 0x07, 0x11)
RED_DARK = RGBColor(0x8B, 0x1A, 0x1A)
RED_DEEP = RGBColor(0x5E, 0x12, 0x12)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x99, 0x99, 0x99)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
GREEN = RGBColor(0x4C, 0xAF, 0x50)
ACCENT_ROSE = RGBColor(0xFF, 0xCC, 0xCC)

# ── Slide dimensions (16:9) ────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]


# ── Helper Functions ───────────────────────────────────

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, italic=False,
             color=WHITE, align=PP_ALIGN.LEFT, font_name="Arial"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_multiline(slide, lines, left, top, width, height,
                  font_size=16, color=WHITE, line_spacing=1.5,
                  bold=False, font_name="Arial"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(line_spacing * 4)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = font_name
    return txBox


def add_left_accent(slide):
    """Rote 6pt vertikale Linie ganz links — Vater-Akzent."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Pt(6), SLIDE_H
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RED
    shape.line.fill.background()


def add_section_label(slide, text, top=Inches(0.35)):
    """Kleines rotes Caps-Label oben — kategorisiert den Slide."""
    add_text(slide, text.upper(), Inches(0.8), top,
             Inches(11), Inches(0.3),
             font_size=11, bold=True, color=RED)


def add_footer(slide, page_num, total=TOTAL_SLIDES):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(6.9), Inches(12.333), Pt(1)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x33)
    shape.line.fill.background()
    add_text(slide, f"{page_num} / {total}",
             Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.4),
             font_size=10, color=GRAY, align=PP_ALIGN.RIGHT)


def add_image_safe(slide, path, left, top, width=None, height=None):
    if os.path.exists(path):
        kwargs = {"image_file": path, "left": left, "top": top}
        if width:
            kwargs["width"] = width
        if height:
            kwargs["height"] = height
        return slide.shapes.add_picture(**kwargs)
    return None


def add_card(slide, left, top, width, height, title, body, bg_color=DARK):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    shape.adjustments[0] = 0.05

    add_text(slide, title, left + Inches(0.3), top + Inches(0.2),
             width - Inches(0.6), Inches(0.5),
             font_size=16, bold=True, color=WHITE)
    add_text(slide, body, left + Inches(0.3), top + Inches(0.65),
             width - Inches(0.6), height - Inches(0.85),
             font_size=12, color=LIGHT_GRAY)


def add_fact_card(slide, left, top, width, height, title, body,
                  accent_color=RED, bg_color=DARK, title_color=WHITE):
    """Karte mit linkem roten Streifen (4pt) — Vater-Stil."""
    # Linker Akzent-Streifen
    strip = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, Pt(4), height
    )
    strip.fill.solid()
    strip.fill.fore_color.rgb = accent_color
    strip.line.fill.background()

    # Main card (ohne runde Ecken, flach)
    card_left = left + Pt(4)
    card_w = width - Pt(4)
    card = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, card_left, top, card_w, height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    card.line.fill.background()

    # Titel
    add_text(slide, title, card_left + Inches(0.25), top + Inches(0.18),
             card_w - Inches(0.4), Inches(0.5),
             font_size=15, bold=True, color=title_color)
    # Body
    add_text(slide, body, card_left + Inches(0.25), top + Inches(0.62),
             card_w - Inches(0.4), height - Inches(0.8),
             font_size=12, color=LIGHT_GRAY)


def add_icon_circle(slide, left, top, size, icon_char, fill_color=RED,
                    icon_color=WHITE, font_size=20):
    """Roter Kreis mit weißem Unicode-Icon — Vater-Anker."""
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = fill_color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = icon_char
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = icon_color
    run.font.name = "Arial"


def chapter_slide(slide, number, title, subtitle):
    set_bg(slide, RED_DARK)
    add_text(slide, f"{number:02d}", Inches(0.8), Inches(1.0),
             Inches(4), Inches(1.5),
             font_size=72, bold=True, color=ACCENT_ROSE,
             font_name="Arial Black")
    add_text(slide, title, Inches(0.8), Inches(2.8),
             Inches(10), Inches(1.5),
             font_size=48, bold=True, color=WHITE, font_name="Arial Black")
    add_text(slide, subtitle, Inches(0.8), Inches(4.5),
             Inches(10), Inches(1),
             font_size=20, italic=False, color=ACCENT_ROSE)


def add_red_badge(slide, text, left, top, width=Inches(2.5), height=Inches(0.5)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RED
    shape.line.fill.background()
    shape.adjustments[0] = 0.3
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Arial"


def add_bullet_list(slide, items, left, top, width, height,
                    icon="✕", icon_color=RED, font_size=14):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(8)

        icon_run = p.add_run()
        icon_run.text = f"{icon}  "
        icon_run.font.size = Pt(font_size)
        icon_run.font.color.rgb = icon_color
        icon_run.font.bold = True
        icon_run.font.name = "Arial"

        text_run = p.add_run()
        text_run.text = item
        text_run.font.size = Pt(font_size)
        text_run.font.color.rgb = LIGHT_GRAY
        text_run.font.name = "Arial"


def add_check_list(slide, items, left, top, width, height, font_size=14):
    add_bullet_list(slide, items, left, top, width, height,
                    icon="✓", icon_color=GREEN, font_size=font_size)


def add_arrow_list(slide, items, left, top, width, height,
                   icon="→", icon_color=RED, font_size=13):
    add_bullet_list(slide, items, left, top, width, height,
                    icon=icon, icon_color=icon_color, font_size=font_size)


# ══════════════════════════════════════════════════════
#  SLIDES — v5 (25 total)
# ══════════════════════════════════════════════════════

# ── SLIDE 1: Titel ─────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, BLACK)
add_left_accent(s)
add_text(s, "Fitness Factory Hattingen",
         Inches(0.8), Inches(1.5), Inches(11), Inches(1.5),
         font_size=52, bold=True, font_name="Arial Black")
add_text(s, "Website-Neuauftritt",
         Inches(0.8), Inches(3.2), Inches(11), Inches(1),
         font_size=36, bold=False, color=RED)
add_text(s, "Von der digitalen Visitenkarte\nzum Mitglieder-Magneten.",
         Inches(0.8), Inches(4.5), Inches(11), Inches(1.2),
         font_size=20, color=GRAY)
add_footer(s, 1)


# ── SLIDE 2: Die eine Frage ────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, BLACK)
add_left_accent(s)
add_text(s, "Die eine Frage.",
         Inches(0.8), Inches(1.2), Inches(11), Inches(1),
         font_size=48, bold=True, font_name="Arial Black")

shape = s.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.8), Inches(2.8), Inches(11.5), Inches(2.5)
)
shape.fill.solid()
shape.fill.fore_color.rgb = DARK
shape.line.fill.background()
shape.adjustments[0] = 0.03

add_text(s, "Hilft deine Website dabei,\nneue Mitglieder zu gewinnen?",
         Inches(1.3), Inches(3.0), Inches(10.5), Inches(1.5),
         font_size=32, bold=True)
add_text(s, "Wenn nein, arbeitet sie gegen dich.",
         Inches(1.3), Inches(4.3), Inches(10.5), Inches(0.6),
         font_size=18, italic=True, color=RED)

add_text(s, "Jedes Mitglied, das den Kursplan nicht findet,\njede Navigation, die verwirrt — kostet dich Geld.",
         Inches(0.8), Inches(5.8), Inches(11), Inches(0.6),
         font_size=14, color=GRAY)
add_text(s, "Die Website ist kein Design-Projekt. Sie ist eine Verkaufsanlage.",
         Inches(0.8), Inches(6.4), Inches(11), Inches(0.4),
         font_size=13, italic=True, color=RED)
add_footer(s, 2)


# ── SLIDE 3: Website heute ─────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, BLACK)
add_left_accent(s)
add_section_label(s, "Status quo")
add_text(s, "Website heute",
         Inches(0.8), Inches(0.7), Inches(6), Inches(0.8),
         font_size=40, bold=True, font_name="Arial Black")
add_text(s, "Was der erste Klick heute kostet",
         Inches(0.8), Inches(1.4), Inches(6), Inches(0.4),
         font_size=14, color=GRAY)

desktop_old = os.path.join(SCREENSHOTS_V4, "desktop-old-top.png")
if not os.path.exists(desktop_old):
    desktop_old = os.path.join(SCREENSHOTS_OLD, "01-startseite.png")
add_image_safe(s, desktop_old, Inches(0.8), Inches(1.9), width=Inches(6.5))

add_text(s, "Live-Screenshot fitness-factory-hattingen.de",
         Inches(0.8), Inches(6.3), Inches(6.5), Inches(0.3),
         font_size=10, color=GRAY, italic=True)

add_text(s, "3 Schwachstellen",
         Inches(7.8), Inches(1.9), Inches(5), Inches(0.5),
         font_size=18, bold=True, color=RED)

add_red_badge(s, "Mobile: Logo unsichtbar, Cookie-Wall vor Hero",
              Inches(7.8), Inches(2.6), width=Inches(4.8), height=Inches(0.55))
add_text(s, "→ Besucher sieht erst Datenschutz-Popup, dann leeren Hero.",
         Inches(7.8), Inches(3.2), Inches(4.8), Inches(0.35),
         font_size=11, color=LIGHT_GRAY)

add_red_badge(s, "One-Pager getarnt als Seitenstruktur",
              Inches(7.8), Inches(3.65), width=Inches(4.8), height=Inches(0.55))
add_text(s, "→ 7 Menüpunkte, davon 4 nur Sprungmarken (#kursplan, #faq, …).",
         Inches(7.8), Inches(4.25), Inches(4.8), Inches(0.35),
         font_size=11, color=LIGHT_GRAY)

add_red_badge(s, "Trust-Signale komplett versteckt",
              Inches(7.8), Inches(4.7), width=Inches(4.8), height=Inches(0.55))
add_text(s, "→ 4,9 ★ bei 167 Google-Bewertungen — auf der Startseite null Hinweis.",
         Inches(7.8), Inches(5.3), Inches(4.8), Inches(0.55),
         font_size=11, color=LIGHT_GRAY)

add_text(s, "Browser-Tab der Live-Seite zeigt nur \"Home –\" — kein Studio-Name, kein Ort.",
         Inches(7.8), Inches(6.0), Inches(4.8), Inches(0.5),
         font_size=11, italic=True, color=GRAY)
add_footer(s, 3)


# ── SLIDE 4: Kapitel 01 — Schwachstellen ───────────────
s = prs.slides.add_slide(blank_layout)
chapter_slide(s, 1, "Schwachstellen",
              "Was Besucher heute erleben —\nund warum sie nicht wiederkommen.")
add_footer(s, 4)


# ── SLIDE 5: NEU — Navigation-Aufgabe (Vater 15) ───────
s = prs.slides.add_slide(blank_layout)
set_bg(s, BLACK)
add_left_accent(s)
add_section_label(s, "Die Aufgabe · Navigation")

# Icon-Circle + Headline
add_icon_circle(s, Inches(0.8), Inches(1.0), Inches(0.8), "☰", font_size=26)
add_text(s, "3 Klicks bis zum Ziel —\nnicht 13 Links zum Rätselraten.",
         Inches(1.85), Inches(0.9), Inches(7.5), Inches(1.8),
         font_size=26, bold=True, font_name="Arial Black")

# Aktuelle Fassung Label
add_text(s, "AKTUELLE FASSUNG",
         Inches(0.8), Inches(3.0), Inches(6), Inches(0.3),
         font_size=11, bold=True, color=RED)

# 5 Problem-Bullets
add_bullet_list(s, [
    "Menüpunkte ohne klare Hierarchie",
    "Conversion-Punkte (Probetraining) gehen unter",
    "Kein Call-to-Action im Header",
    "Mobil: unsortierte Textliste ohne Icons",
    "»Info«-Menüpunkt ohne erkennbare Funktion",
], Inches(0.8), Inches(3.35), Inches(6.5), Inches(2.8), font_size=14)

# Rechts: kleiner Screenshot der alten Nav-Leiste mit "Kommentar"-Box
add_image_safe(s, os.path.join(SCREENSHOTS_V4, "desktop-old-top.png"),
               Inches(7.8), Inches(2.5), width=Inches(4.7))
# Kommentar-Overlay (Post-it-Stil)
comment = s.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(8.0), Inches(4.5), Inches(4.5), Inches(1.1)
)
comment.fill.solid()
comment.fill.fore_color.rgb = DARK
comment.line.color.rgb = RED
comment.line.width = Pt(1.5)
add_text(s, "Menü Desktop",
         Inches(8.15), Inches(4.6), Inches(4.3), Inches(0.35),
         font_size=12, bold=True, color=RED)
add_text(s, "13 Einträge ohne Priorisierung.\nProbetraining + Telefon verschwinden in der Masse.",
         Inches(8.15), Inches(4.95), Inches(4.3), Inches(0.65),
         font_size=10, color=LIGHT_GRAY)

add_text(s, "→ Besucher scrollen weg, bevor sie klicken.",
         Inches(0.8), Inches(6.5), Inches(11), Inches(0.3),
         font_size=13, italic=True, color=RED)
add_footer(s, 5)


# ── SLIDE 6: Navigation & Header (Vorher/Nachher, war v4 Slide 5) ──
s = prs.slides.add_slide(blank_layout)
set_bg(s, BLACK)
add_left_accent(s)
add_section_label(s, "Aktuelle Fassung vs. Relaunch")
add_text(s, "Navigation & Header",
         Inches(0.8), Inches(0.7), Inches(11), Inches(0.8),
         font_size=36, bold=True, font_name="Arial Black")

add_image_safe(s, os.path.join(SCREENSHOTS_V4, "desktop-old-top.png"),
               Inches(0.8), Inches(1.7), width=Inches(5.7))
add_image_safe(s, os.path.join(SCREENSHOTS_V4, "desktop-new-top.png"),
               Inches(6.8), Inches(1.7), width=Inches(5.7))

add_text(s, "VORHER (Live)", Inches(0.8), Inches(5.2), Inches(3), Inches(0.4),
         font_size=12, bold=True, color=RED)
add_bullet_list(s, [
    "13 Links ohne Hierarchie",
    "Kein CTA im Header",
    "Cookie-Wall blockiert Hero",
    "Keine Telefon-/WhatsApp-Shortcut"
], Inches(0.8), Inches(5.6), Inches(5.7), Inches(1.2), font_size=12)

add_text(s, "NACHHER (Relaunch)", Inches(6.8), Inches(5.2), Inches(4), Inches(0.4),
         font_size=12, bold=True, color=GREEN)
add_check_list(s, [
    "6 klare Hauptpunkte + Logo sichtbar",
    "\"Probetraining buchen\" + \"Jetzt anrufen\"",
    "Hero above-the-fold, keine Popups",
    "Sticky CTA beim Scrollen"
], Inches(6.8), Inches(5.6), Inches(5.7), Inches(1.2), font_size=12)

add_footer(s, 6)


# ── SLIDE 7: NEU — Kursplan-Aufgabe (Vater 5) ──────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, BLACK)
add_left_accent(s)
add_section_label(s, "Die Aufgabe · Kursplan")

add_icon_circle(s, Inches(0.8), Inches(1.0), Inches(0.8), "◉", font_size=24)
add_text(s, "Der Kursplan ist eines der wichtigsten\nWerkzeuge deiner Mitglieder.",
         Inches(1.85), Inches(0.9), Inches(10.5), Inches(1.8),
         font_size=26, bold=True, font_name="Arial Black")

add_text(s, "Er entscheidet, ob ein Mitglied heute Abend zum Kurs kommt —\noder auf der Couch bleibt.",
         Inches(0.8), Inches(2.8), Inches(11.5), Inches(0.8),
         font_size=15, color=GRAY)

# 4 Check-Kriterien
criteria = [
    ("Sofort sichtbar", "kein Scrollen, kein Suchen"),
    ("Perfekt lesbar", "auf jedem Gerät, in jeder Schriftgröße"),
    ("Immer aktuell", "Änderungen und Ausfälle sofort sichtbar"),
    ("Motivierend", "mit Trainer-Info, Beschreibung, Anmelde-Option"),
]
for i, (title, desc) in enumerate(criteria):
    y = Inches(3.9) + Inches(0.7) * i
    add_icon_circle(s, Inches(0.9), y, Inches(0.4), "✓",
                    fill_color=GREEN, font_size=14)
    add_text(s, title, Inches(1.5), y + Inches(0.02),
             Inches(3.5), Inches(0.4),
             font_size=15, bold=True)
    add_text(s, "— " + desc, Inches(5.0), y + Inches(0.02),
             Inches(7.5), Inches(0.4),
             font_size=13, color=LIGHT_GRAY)

add_footer(s, 7)


# ── SLIDE 8: Kursplan & Bewertungen (v4 Slide 6) ───────
s = prs.slides.add_slide(blank_layout)
set_bg(s, BLACK)
add_left_accent(s)
add_section_label(s, "Aktuelle Fassung vs. Relaunch")
add_text(s, "Kursplan & Bewertungen",
         Inches(0.8), Inches(0.7), Inches(11), Inches(0.8),
         font_size=36, bold=True, font_name="Arial Black")

add_fact_card(s, Inches(0.8), Inches(1.8), Inches(5.85), Inches(2.25),
              "Kursplan — vorher",
              "Reine Tabelle ohne Filter.\nKeine Kursbeschreibungen.\nKeine Anmelde-Option.\n16 Kurse — keiner wird emotional verkauft.",
              accent_color=RED)
add_fact_card(s, Inches(6.85), Inches(1.8), Inches(5.85), Inches(2.25),
              "Kursplan — nachher",
              "Klickbares Grid mit Trainer-Foto.\nFilter: Cardio · Kraft · Entspannung.\nJeder Kurs mit Beschreibung + CTA.\n→ Aus Tabelle wird Kurs-Erlebnis.",
              accent_color=GREEN)
add_fact_card(s, Inches(0.8), Inches(4.25), Inches(5.85), Inches(2.25),
              "4,9 ★ bei 167 Bewertungen — versteckt",
              "Google-Bewertungen nirgendwo\nauf der Startseite sichtbar.\nKein Link zum Google-Profil.\n→ Stärkster Beweis komplett ungenutzt.",
              accent_color=RED)
add_fact_card(s, Inches(6.85), Inches(4.25), Inches(5.85), Inches(2.25),
              "4,9 ★ — jetzt auf jeder Seite",
              "Trust-Bar mit echten Stimmen.\nGoogle-Logo + Sterne-Widget.\nCounter-Animation beim Scrollen.\n→ 167 stille Bewertungen werden Verkäufer.",
              accent_color=GREEN)
add_footer(s, 8)


# ── SLIDE 9: Details die Vertrauen kosten (v4 Slide 7) ──
s = prs.slides.add_slide(blank_layout)
set_bg(s, BLACK)
add_left_accent(s)
add_section_label(s, "Was Besucher auf der Live-Site sofort sehen")
add_text(s, "Details, die Vertrauen kosten.",
         Inches(0.8), Inches(0.7), Inches(11), Inches(0.7),
         font_size=32, bold=True, font_name="Arial Black")

items = [
    ("Mobile: Logo unsichtbar", "Auf 68 % der Zugriffe fehlt die Marke im Header — kein Wiedererkennungswert."),
    ("Kursplan nicht aktuell", "Angezeigte Kurse stimmen nicht mit Realität überein — sofortiger Vertrauensbruch."),
    ("Anker-Nav wirkt wie Seiten", "7 Menüpunkte, 4 davon springen nur zu Abschnitten — Besucher verliert Orientierung."),
    ("Keine Route-/Anfahrts-Funktion", "Keine eingebettete Karte, kein „Route planen\"-Button — Weg zum Studio fehlt."),
    ("SEPA-Formular öffentlich", "Bankdaten-Seite ungeschützt im Web — DSGVO- und Sicherheitsrisiko."),
    ("GMX-Mailadresse", "Wirkt hobbyhaft neben Premium-Preisen — Domain-Mail ist Standard."),
    ("Kein Team- / Trainer-Bereich", "Kein Gesicht, kein Name — „familiär\" wird behauptet, nicht gezeigt."),
]

row_h = Inches(0.65)
start_y = Inches(1.65)
for i, (title, desc) in enumerate(items):
    y = start_y + row_h * i
    add_text(s, "⚠", Inches(0.8), y, Inches(0.45), Inches(0.4),
             font_size=14, color=RED)
    add_text(s, title, Inches(1.3), y, Inches(3.8), Inches(0.4),
             font_size=13, bold=True)
    add_text(s, desc, Inches(5.2), y, Inches(7.5), Inches(0.5),
             font_size=11, color=LIGHT_GRAY)

add_text(s, "→ Das Studio ist besser als sein Auftritt. Genau das kostet Mitglieder.",
         Inches(0.8), Inches(6.5), Inches(11), Inches(0.4),
         font_size=12, italic=True, color=GRAY)
add_footer(s, 9)


# ── SLIDE 10: Kapitel 02 — Das Logo ────────────────────
s = prs.slides.add_slide(blank_layout)
chapter_slide(s, 2, "Das Logo",
              "Die Marke muss auf den ersten Blick\nals Fitnessstudio erkennbar sein.")
add_footer(s, 10)


# ── SLIDE 11: NEU — Logo-Aufgabe (Vater 10) ────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, BLACK)
add_left_accent(s)
add_section_label(s, "Die Aufgabe · Logo")
add_text(s, "Ein Logo muss in 1 Sekunde sagen:\n»Hier trainierst du.«",
         Inches(0.8), Inches(0.75), Inches(7.5), Inches(2.2),
         font_size=28, bold=True, font_name="Arial Black")

# AKTUELLE FASSUNG Label
add_text(s, "AKTUELLE FASSUNG",
         Inches(0.8), Inches(3.2), Inches(5), Inches(0.3),
         font_size=11, bold=True, color=RED)

# 4 Probleme als X-Liste mit Titel+Body
problems = [
    ("Man liest »ITNESS« statt »FITNESS«",
     "Das verfremdete F verschwindet im Blitz-Symbol — übrig bleibt »ITNESS FACTORY«."),
    ("IT statt Fitness?",
     "»ITNESS FACTORY« klingt nach IT-Dienstleister, Programmierung oder Computer-Firma."),
    ("Der Blitz hat keine Branchenzuordnung",
     "Elektrobetrieb, eSports, Auto-Tuning, Energiedrink — alles, nur nicht Fitness."),
    ("Keine sofortige Fitness-Assoziation",
     "Nichts am Symbol kommuniziert Sport, Kraft oder Bewegung."),
]
for i, (title, body) in enumerate(problems):
    y = Inches(3.55) + Inches(0.72) * i
    add_text(s, "✕", Inches(0.85), y, Inches(0.4), Inches(0.4),
             font_size=15, bold=True, color=RED)
    add_text(s, title, Inches(1.35), y - Inches(0.02),
             Inches(6.5), Inches(0.35),
             font_size=13, bold=True)
    add_text(s, body, Inches(1.35), y + Inches(0.28),
             Inches(6.5), Inches(0.4),
             font_size=11, color=LIGHT_GRAY)

# Logo-Bild rechts
add_image_safe(s, LOGO_PATH, Inches(8.6), Inches(2.0), width=Inches(4.0))

# Fazit unten rechts
add_text(s, "→ Wer das Logo sieht, liest »ITNESS« —",
         Inches(8.3), Inches(5.8), Inches(4.8), Inches(0.35),
         font_size=12, italic=True, bold=True, color=RED)
add_text(s, "und denkt an alles, nur nicht an ein Fitnessstudio.",
         Inches(8.3), Inches(6.15), Inches(4.8), Inches(0.35),
         font_size=12, italic=True, color=RED)
add_footer(s, 11)


# ── SLIDE 12: NEU — Logo-Tiefenanalyse (Vater 11) ──────
s = prs.slides.add_slide(blank_layout)
set_bg(s, BLACK)
add_left_accent(s)
add_section_label(s, "Logo-Tiefenanalyse")

add_icon_circle(s, Inches(0.8), Inches(1.0), Inches(0.8), "⌕", font_size=26)
add_text(s, "Warum der Blitz grundsätzlich\nhinterfragt werden muss.",
         Inches(1.85), Inches(0.9), Inches(10.5), Inches(1.8),
         font_size=26, bold=True, font_name="Arial Black")

# 2x2 Grid Problem-Karten mit linkem roten Streifen
card_problems = [
    ("3D-Metallic-Effekt",
     "Chrome-Verläufe wirken nach 2005.\nModerne Fitness-Marken setzen\nauf klares, flaches Design —\nFitX, McFIT, John Reed, Urban\nSports Club machen es vor."),
    ("»ITNESS« statt »FITNESS«",
     "Das verfremdete F verschwindet —\nübrig bleibt ITNESS FACTORY.\nDas klingt nach IT-Firma, nicht\nnach Fitnessstudio. Drei Schrift-\nschnitte verstärken das Chaos."),
    ("Zielgruppen-Signal",
     "Rot-Metallic auf Schwarz wirkt\naggressiv und maskulin. Passt\nnicht zu »familiäres Studio für\nalle« — schließt Frauen, Ältere\nund Anfänger optisch aus."),
    ("Keine Skalierbarkeit",
     "Feine Metallic-Details gehen bei\nkleinen Größen verloren: Favicon,\nInstagram-Profilbild, Stickerei,\nStempel — überall unleserlich."),
]
card_w = Inches(5.9)
card_h = Inches(1.95)
positions = [
    (Inches(0.8), Inches(2.75)),
    (Inches(6.8), Inches(2.75)),
    (Inches(0.8), Inches(4.8)),
    (Inches(6.8), Inches(4.8)),
]
for (x, y), (title, body) in zip(positions, card_problems):
    add_fact_card(s, x, y, card_w, card_h, title, body, accent_color=RED)

add_footer(s, 12)


# ── SLIDE 13: NEU — Logo Option A vs B (Vater 12) ──────
s = prs.slides.add_slide(blank_layout)
set_bg(s, BLACK)
add_left_accent(s)
add_section_label(s, "Der Blitz — bleiben oder gehen?")

add_icon_circle(s, Inches(0.8), Inches(0.95), Inches(0.8), "⚡", font_size=26)
add_text(s, "Ein Symbol muss die Marke stärken —\nnicht erklärungsbedürftig sein.",
         Inches(1.85), Inches(0.85), Inches(10.5), Inches(1.8),
         font_size=24, bold=True, font_name="Arial Black")

# Option A (links) — grau-betont, Blitz bleibt
opt_w = Inches(5.95)
opt_h = Inches(3.7)
opt_y = Inches(2.9)

# Option A Header (grau)
a_header = s.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.8), opt_y, opt_w, Inches(0.65)
)
a_header.fill.solid()
a_header.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x33)
a_header.line.fill.background()
add_text(s, "⚡  Option A: Blitz bleibt",
         Inches(1.0), opt_y + Inches(0.13), opt_w - Inches(0.4), Inches(0.5),
         font_size=16, bold=True, color=LIGHT_GRAY)

# Option A Body
a_body = s.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0.8), opt_y + Inches(0.65), opt_w, opt_h - Inches(0.65)
)
a_body.fill.solid()
a_body.fill.fore_color.rgb = DARK
a_body.line.fill.background()

add_text(s, "Dann muss er:",
         Inches(1.05), opt_y + Inches(0.85),
         opt_w - Inches(0.4), Inches(0.4),
         font_size=14, bold=True)

a_bullets = [
    "als »F« eindeutig lesbar werden — nie wieder »ITNESS« statt »FITNESS«",
    "vereinfacht werden (kein 3D, kein Metallic)",
    "mit einem Fitness-Signal ergänzt werden (Hantel-Silhouette, Bewegungs-Swoosh)",
    "in jeder Größe funktionieren (Favicon bis Fassade)",
]
add_arrow_list(s, a_bullets,
               Inches(1.05), opt_y + Inches(1.35),
               opt_w - Inches(0.4), opt_h - Inches(1.5),
               icon="→", icon_color=LIGHT_GRAY, font_size=12)

# Option B (rechts) — rot-betont, Neues Symbol
b_x = Inches(6.9)
b_header = s.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, b_x, opt_y, opt_w, Inches(0.65)
)
b_header.fill.solid()
b_header.fill.fore_color.rgb = RED_DARK
b_header.line.fill.background()
add_text(s, "⚡  Option B: Neues Symbol",
         b_x + Inches(0.2), opt_y + Inches(0.13),
         opt_w - Inches(0.4), Inches(0.5),
         font_size=16, bold=True, color=WHITE)

b_body = s.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, b_x, opt_y + Inches(0.65), opt_w, opt_h - Inches(0.65)
)
b_body.fill.solid()
b_body.fill.fore_color.rgb = DARK
b_body.line.fill.background()

add_text(s, "Ein neues Zeichen, das:",
         b_x + Inches(0.25), opt_y + Inches(0.85),
         opt_w - Inches(0.4), Inches(0.4),
         font_size=14, bold=True)

b_bullets = [
    "sofort »Fitness« kommuniziert",
    "als Wortmarke + Icon funktioniert",
    "den Namen »Factory« visuell aufgreift (industrielle Typografie, Zahnrad-Element)",
    "modern, flach und skalierbar ist",
]
add_arrow_list(s, b_bullets,
               b_x + Inches(0.25), opt_y + Inches(1.35),
               opt_w - Inches(0.4), opt_h - Inches(1.5),
               icon="→", icon_color=ACCENT_ROSE, font_size=12)

add_text(s, "Empfehlung: Option B — sauberer Neustart statt halbherziger Reparatur.",
         Inches(0.8), Inches(6.55), Inches(11.5), Inches(0.3),
         font_size=12, italic=True, bold=True, color=RED, align=PP_ALIGN.CENTER)
add_footer(s, 13)


# ── SLIDE 14: Kapitel 03 — Der Neuauftritt ─────────────
s = prs.slides.add_slide(blank_layout)
chapter_slide(s, 3, "Der Neuauftritt",
              "Gebaut für Google.\nGebaut für Handys.\nGebaut für Kunden.")
add_footer(s, 14)


# ── SLIDE 15: Header & Hero (v4 Slide 9) ───────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, BLACK)
add_left_accent(s)
add_section_label(s, "Ideale Umsetzung · Header & Hero")
add_text(s, "Header & Hero",
         Inches(0.8), Inches(0.7), Inches(11), Inches(0.8),
         font_size=36, bold=True, font_name="Arial Black")
add_text(s, "Die ersten 3 Sekunden verkaufen.\nDer Rest der Seite bestätigt nur noch.",
         Inches(0.8), Inches(1.5), Inches(5), Inches(0.8),
         font_size=16, color=GRAY)

add_image_safe(s, os.path.join(SCREENSHOTS_NEW, "01-hero.png"),
               Inches(6.5), Inches(0.7), width=Inches(6.3))

add_check_list(s, [
    "Starke Headline: \"Dein Gym in Hattingen\"",
    "TrustBar: 4,9 ★ · 167 Bewertungen · 16 Kurse",
    "Primär-CTA: \"Probetraining sichern\" (Rot)",
    "Sekundär: \"Jetzt anrufen\" (Outline)",
    "Sticky CTA-Bar nach 30% Scroll"
], Inches(0.8), Inches(2.6), Inches(5.5), Inches(3.5), font_size=14)

add_text(s, "→ Claim, Beweis, CTA — alles above the fold.",
         Inches(0.8), Inches(6.5), Inches(11), Inches(0.4),
         font_size=13, italic=True, color=GRAY)
add_footer(s, 15)


# ── SLIDE 16: Leistungen & Preise (v4 Slide 10) ────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, BLACK)
add_left_accent(s)
add_section_label(s, "Ideale Umsetzung · Leistungen & Preise")
add_text(s, "Leistungen & Preise",
         Inches(0.8), Inches(0.7), Inches(11), Inches(0.8),
         font_size=36, bold=True, font_name="Arial Black")

add_card(s, Inches(0.8), Inches(1.7), Inches(5.85), Inches(4.6),
         "Leistungen — neu",
         "Aus generischen Icons werden\nechte Räume und echte Versprechen.\n\n"
         "• Jede Leistung mit echtem Foto\n"
         "• Nutzen-Claims statt Feature-Listen\n"
         "• 8 Kernleistungen im 2×4 Grid\n"
         "• CTA: \"Alle Leistungen live erleben\"")

add_card(s, Inches(6.85), Inches(1.7), Inches(5.85), Inches(4.6),
         "Preise — neu",
         "3 Pakete. 1 klarer Favorit.\nPreise werden verkauft, nicht gelistet.\n\n"
         "• Basis (55€) · Standard (45€) · Premium (35€)\n"
         "• \"Beliebteste Wahl\" rot markiert\n"
         "• Fremdgeh-Aktion als Trigger\n"
         "• CTA pro Tarif: \"Jetzt Probetraining\"")

add_text(s, "→ \"Alles inklusive\" wird \"Das bekommst du bei uns — und so fühlt es sich an.\"",
         Inches(0.8), Inches(6.5), Inches(11), Inches(0.4),
         font_size=13, italic=True, color=GRAY)
add_footer(s, 16)


# ── SLIDE 17: Kursplan & FAQ (v4 Slide 11) ─────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, BLACK)
add_left_accent(s)
add_section_label(s, "Ideale Umsetzung · Kursplan & FAQ")
add_text(s, "Kursplan & FAQ",
         Inches(0.8), Inches(0.7), Inches(11), Inches(0.8),
         font_size=36, bold=True, font_name="Arial Black")

add_image_safe(s, os.path.join(SCREENSHOTS_NEW, "03-kursplan.png"),
               Inches(0.8), Inches(1.7), width=Inches(5.7))

add_text(s, "FAQ — strukturiert",
         Inches(7.0), Inches(1.7), Inches(5.5), Inches(0.5),
         font_size=20, bold=True)
add_check_list(s, [
    "3 Kategorien: Mitgliedschaft · Kurse · Wellness",
    "Ausklappbar, max. 3 Sätze pro Antwort",
    "Kein Juristen-Deutsch",
    "Abschluss-CTA: \"Jetzt WhatsApp schreiben\"",
    "Versprechen: Antwort in unter 2 Stunden"
], Inches(7.0), Inches(2.4), Inches(5.5), Inches(3.5), font_size=13)

add_text(s, "→ Jede beantwortete Frage ist eine gewonnene Mitgliedschaft.",
         Inches(0.8), Inches(6.5), Inches(11), Inches(0.4),
         font_size=13, italic=True, color=GRAY)
add_footer(s, 17)


# ── SLIDE 18: Kontakt & Footer (v4 Slide 12) ───────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, BLACK)
add_left_accent(s)
add_section_label(s, "Ideale Umsetzung · Kontakt & Footer")
add_text(s, "Kontakt & Footer",
         Inches(0.8), Inches(0.7), Inches(11), Inches(0.8),
         font_size=36, bold=True, font_name="Arial Black")

add_image_safe(s, os.path.join(SCREENSHOTS_NEW, "05-kontakt.png"),
               Inches(0.8), Inches(1.7), width=Inches(5.7))

add_check_list(s, [
    "Google-Maps-Embed mit echtem Standort",
    "\"Route planen\"-Button (iOS / Android)",
    "Öffnungszeiten live: \"Jetzt geöffnet · schließt um 22:00\"",
    "WhatsApp-Shortcut direkt daneben",
    "Footer als zweite Startseite mit allen Infos"
], Inches(7.0), Inches(2.0), Inches(5.5), Inches(3.5), font_size=14)

add_text(s, "→ Aus \"Weitere Fragen?\" wird \"Ich komme jetzt vorbei.\"",
         Inches(0.8), Inches(6.5), Inches(11), Inches(0.4),
         font_size=13, italic=True, color=GRAY)
add_footer(s, 18)


# ── SLIDE 19: Mobile-First (v4 Slide 13) ───────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, BLACK)
add_left_accent(s)
add_section_label(s, "Aktuelle Fassung vs. Relaunch · Mobil")
add_text(s, "Mobile-First",
         Inches(0.8), Inches(0.7), Inches(11), Inches(0.7),
         font_size=36, bold=True, font_name="Arial Black")
add_text(s, "68 % der Besucher kommen mobil. Ihr erster Eindruck entscheidet alles.",
         Inches(0.8), Inches(1.45), Inches(11), Inches(0.4),
         font_size=15, color=GRAY)

mob_h = Inches(4.1)
add_image_safe(s, os.path.join(SCREENSHOTS_V4, "mobile-old-top.png"),
               Inches(2.3), Inches(2.0), height=mob_h)
add_image_safe(s, os.path.join(SCREENSHOTS_V4, "mobile-new-top.png"),
               Inches(8.5), Inches(2.0), height=mob_h)

add_text(s, "VORHER", Inches(1.8), Inches(6.2), Inches(3.5), Inches(0.3),
         font_size=12, bold=True, color=RED, align=PP_ALIGN.CENTER)
add_text(s, "Cookie-Wall, kein Logo, kein CTA",
         Inches(1.8), Inches(6.48), Inches(3.5), Inches(0.3),
         font_size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_text(s, "NACHHER", Inches(8.0), Inches(6.2), Inches(3.5), Inches(0.3),
         font_size=12, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
add_text(s, "Logo · Hero · Probetraining · Anruf · WhatsApp",
         Inches(8.0), Inches(6.48), Inches(3.5), Inches(0.3),
         font_size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_footer(s, 19)


# ── SLIDE 20: Kapitel 04 — Mehr als eine Website ───────
s = prs.slides.add_slide(blank_layout)
chapter_slide(s, 4, "Mehr als eine Website",
              "Was zusätzlich im Paket steckt.")
add_footer(s, 20)


# ── SLIDE 21: Dashboard (v4 Slide 15) ──────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, BLACK)
add_left_accent(s)
add_section_label(s, "Bonus · Content-Dashboard")
add_text(s, "Dashboard",
         Inches(0.8), Inches(0.7), Inches(11), Inches(0.8),
         font_size=36, bold=True, font_name="Arial Black")
add_text(s, "Texte, Kurse, Öffnungszeiten — alles selbst ändern.\nOhne Webentwickler. Ohne Wartezeit.",
         Inches(0.8), Inches(1.5), Inches(11), Inches(0.8),
         font_size=18, color=GRAY)

add_image_safe(s, os.path.join(FF_SLIDES, "FitnessFactory.030.jpeg"),
               Inches(0.8), Inches(2.6), width=Inches(11.5))

add_text(s, "→ Die Website wird zur Plattform. Das Team behält die Kontrolle.",
         Inches(0.8), Inches(6.5), Inches(11), Inches(0.4),
         font_size=13, italic=True, color=GRAY)
add_footer(s, 21)


# ── SLIDE 22: Schilder vorher (v4 Slide 16) ────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, BLACK)
add_left_accent(s)
add_section_label(s, "Bonus · Schilder vor Ort")
add_text(s, "Nicht nur digital. Auch vor Ort.",
         Inches(0.8), Inches(0.7), Inches(12), Inches(0.8),
         font_size=32, bold=True, font_name="Arial Black")
add_text(s, "Neue Geräte, neue Duschen, neuer Boden — Hinweise aber noch mit Klebeband.",
         Inches(0.8), Inches(1.45), Inches(12), Inches(0.4),
         font_size=15, color=GRAY)

img_w = Inches(3.8)
img_h = Inches(3.7)
positions_img = [Inches(0.55), Inches(4.75), Inches(8.95)]
images_before = [
    (os.path.join(SCREENSHOTS_V4, "brustpresse_voher.jpg"),
     "Brustpresse", "Loser Zettel in einer Folie"),
    (os.path.join(SCREENSHOTS_V4, "sauna_voher.jpg"),
     "Sauna", "Papier auf Metallschild geklebt"),
    (os.path.join(SCREENSHOTS_V4, "spinde_voher.jpg"),
     "Spinde", "Hinweise handgeklebt"),
]
for x, (path, label, hint) in zip(positions_img, images_before):
    add_image_safe(s, path, x, Inches(1.95), width=img_w, height=img_h)
    add_text(s, label, x, Inches(5.75), img_w, Inches(0.35),
             font_size=16, bold=True, align=PP_ALIGN.CENTER, color=RED)
    add_text(s, hint, x, Inches(6.1), img_w, Inches(0.3),
             font_size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_text(s, "→ Premium-Preis verlangt Premium-Auftreten — auch an der Wand.",
         Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.3),
         font_size=12, italic=True, bold=True, color=RED, align=PP_ALIGN.CENTER)
add_footer(s, 22)


# ── SLIDE 23: Schilder-System im CI-Design (v4 Slide 17) ──
s = prs.slides.add_slide(blank_layout)
set_bg(s, BLACK)
add_left_accent(s)
add_section_label(s, "Ideale Umsetzung · Schilder-System")
add_text(s, "Ein Schilder-System im CI-Design.",
         Inches(0.8), Inches(0.7), Inches(12), Inches(0.8),
         font_size=32, bold=True, font_name="Arial Black")
add_text(s, "Konsistent. Markenkonform. Produktionsreif — für heutige und zukünftige Hinweise.",
         Inches(0.8), Inches(1.45), Inches(12), Inches(0.4),
         font_size=15, color=GRAY)

images_after = [
    (os.path.join(SCHILDER, "Brustpresse", "Brustpresse_nachher.jpeg"),
     "Brustpresse", "Geräte-Hinweis im Rahmen"),
    (os.path.join(SCHILDER, "Sauna", "Sauna_nachher.jpeg"),
     "Sauna", "Regeln sauber im Rahmen"),
    (os.path.join(SCHILDER, "Spinde", "Spinde_nacher.jpeg"),
     "Spinde", "CI-Schild mit Logo"),
]
for x, (path, label, hint) in zip(positions_img, images_after):
    add_image_safe(s, path, x, Inches(1.95), width=img_w, height=img_h)
    add_text(s, label, x, Inches(5.75), img_w, Inches(0.35),
             font_size=16, bold=True, align=PP_ALIGN.CENTER, color=GREEN)
    add_text(s, hint, x, Inches(6.1), img_w, Inches(0.3),
             font_size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_text(s, "→ Ein Logo, eine Website, ein Studio — alles aus einem Guss.",
         Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.3),
         font_size=12, italic=True, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
add_footer(s, 23)


# ── SLIDE 24: Der Preis des Nichtstuns (v4 Slide 18) ───
s = prs.slides.add_slide(blank_layout)
set_bg(s, BLACK)
add_left_accent(s)
add_section_label(s, "Zusammenfassung")
add_text(s, "Der Preis des Nichtstuns",
         Inches(0.8), Inches(0.7), Inches(11), Inches(0.8),
         font_size=36, bold=True, font_name="Arial Black")
add_text(s, "Was eine halbfertige Website jeden Monat kostet.",
         Inches(0.8), Inches(1.5), Inches(11), Inches(0.4),
         font_size=14, color=GRAY)

cards = [
    ("68 %", "kommen mobil", "Und stolpern über\nCookie-Wall + fehlendes Logo."),
    ("3 Sek.", "Entscheidungszeit", "Dann ist der\nBesucher weg."),
    ("2.430 €", "entgangener Umsatz/Monat", "~90 verlorene Kontakte,\n~5 % davon = Mitglieder."),
    ("3 × 0 ROI", "in 4 Jahren", "Websites gekauft — ohne\nVerkaufs-Wirkung."),
]
card_w_18 = Inches(2.85)
gap = Inches(0.15)
total_w = card_w_18 * 4 + gap * 3
start_x = (SLIDE_W - total_w) / 2
for i, (number, label, detail) in enumerate(cards):
    x = start_x + (card_w_18 + gap) * i
    card = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x, Inches(2.2), card_w_18, Inches(4.0)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RED_DARK if i == 3 else DARK
    card.line.fill.background()
    card.adjustments[0] = 0.05

    n_size = 34 if len(number) > 6 else 44
    add_text(s, number, x + Inches(0.2), Inches(2.5),
             card_w_18 - Inches(0.4), Inches(1.2),
             font_size=n_size, bold=True,
             color=ACCENT_ROSE if i == 3 else RED,
             font_name="Arial Black")
    add_text(s, label, x + Inches(0.2), Inches(3.9),
             card_w_18 - Inches(0.4), Inches(0.5),
             font_size=13, bold=True)
    add_text(s, detail, x + Inches(0.2), Inches(4.6),
             card_w_18 - Inches(0.4), Inches(1.4),
             font_size=11, color=LIGHT_GRAY if i == 3 else GRAY)

add_text(s, "Jede Woche ohne verkaufende Website kostet Mitgliedschaften.",
         Inches(0.8), Inches(6.45), Inches(11), Inches(0.4),
         font_size=15, italic=True, color=RED, align=PP_ALIGN.CENTER)
add_footer(s, 24)


# ── SLIDE 25: Nächste Schritte + Abschluss (fusioniert) ──
s = prs.slides.add_slide(blank_layout)
set_bg(s, BLACK)
add_left_accent(s)
add_section_label(s, "Nächste Schritte")
add_text(s, "Die Website ist fertig.",
         Inches(0.8), Inches(0.7), Inches(11), Inches(0.9),
         font_size=40, bold=True, font_name="Arial Black")
add_text(s, "Sie wartet auf dein Go — und in 4 Wochen ist sie live.",
         Inches(0.8), Inches(1.6), Inches(11), Inches(0.4),
         font_size=16, color=GRAY)

steps = [
    ("1", "Freigabe Farbwelt & Typografie", "Woche 1"),
    ("2", "Logo-Neugestaltung & Bildbearbeitung", "Woche 1–2"),
    ("3", "Website-Relaunch mit neuem Kursplan", "Woche 2–3"),
    ("4", "Instagram-Update & Highlight-Covers", "Woche 3–4"),
    ("5", "Live-Gang & Feinschliff", "Woche 4"),
]

for i, (num, title, week) in enumerate(steps):
    y = Inches(2.3) + Inches(0.82) * i
    circle = s.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(0.8), y, Inches(0.55), Inches(0.55)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = RED
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = num
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Arial"

    add_text(s, title, Inches(1.6), y + Inches(0.02),
             Inches(6), Inches(0.35),
             font_size=16, bold=True)
    add_text(s, week, Inches(1.6), y + Inches(0.35),
             Inches(6), Inches(0.3),
             font_size=12, color=GRAY)

bereit = s.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(8.5), Inches(2.6), Inches(4), Inches(3.6)
)
bereit.fill.solid()
bereit.fill.fore_color.rgb = RED_DARK
bereit.line.fill.background()
bereit.adjustments[0] = 0.05

add_text(s, "Bereit?", Inches(9.0), Inches(2.85),
         Inches(3), Inches(0.8),
         font_size=36, bold=True, font_name="Arial Black")
add_text(s, "Nach Freigabe:\n7 Tage bis Live-Gang.",
         Inches(9.0), Inches(3.75), Inches(3), Inches(1),
         font_size=15, color=ACCENT_ROSE)
add_text(s, "Lass uns starten.",
         Inches(9.0), Inches(4.95), Inches(3), Inches(0.4),
         font_size=18, bold=True, italic=True, color=WHITE)
add_text(s, "Diesmal gebaut,\num zu bleiben.",
         Inches(9.0), Inches(5.45), Inches(3), Inches(0.8),
         font_size=13, italic=True, color=ACCENT_ROSE)

add_text(s, "Luke Kozik  ·  Web Development & Digital Strategy",
         Inches(0.8), Inches(6.55), Inches(7), Inches(0.3),
         font_size=11, color=GRAY)
add_footer(s, 25)


# ── Save ───────────────────────────────────────────────
prs.save(OUTPUT)
print(f"✓ Präsentation erstellt: {OUTPUT}")
print(f"  → {len(prs.slides)} Slides")
