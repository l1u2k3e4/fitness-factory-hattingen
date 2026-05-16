#!/usr/bin/env python3
"""
Fitness Factory Hattingen — Relaunch-Präsentation v4
21 Slides im Karteikarten-Stil (PowerPoint)

Änderungen gegenüber v3:
 - AkquiseFlow-Branding entfernt (Logo-Platzhalter)
 - Slide 3: ehrliche Fakten statt falsche SEO/404-Claims, + Conversion-Impact-Pfeile
 - Slide 5 + 13: echte Website-Screenshots statt v1-Text-Bilder
 - Slide 7: Tippfehler-Punkt raus, 4 neue Schwachpunkte rein (7 statt 6)
 - NEU Slide 16 + 17: Studio-Auftreten (Klebeband-Zettel vorher / CI-Schilder nachher)
 - Instagram-Slide (v3 Slide 16) entfernt
 - Preis des Nichtstuns: 4. Karte "3 Websites in 4 Jahren = 0 ROI"
 - Überzeugungs-Framing geschärft (Alex hat 3 Websites in 4 Jahren gekauft)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Paths ──────────────────────────────────────────────
BASE = os.path.dirname(os.path.abspath(__file__))
SCREENSHOTS_NEW = os.path.join(BASE, "figma-screenshots", "new")
SCREENSHOTS_OLD = os.path.join(BASE, "figma-screenshots", "old")
SCREENSHOTS_V4 = os.path.join(BASE, "figma-screenshots", "v4")
FF_SLIDES = os.path.join(BASE, "FitnessFactory")
SCHILDER = "/Users/lukekozik/Documents/Programme/n8n/Jobs/FitnessFactory/Schilder"
OUTPUT = os.path.join(BASE, "FitnessFactory_Relaunch_v4.pptx")

# ── Brand Colors ───────────────────────────────────────
BLACK = RGBColor(0x11, 0x11, 0x11)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
RED = RGBColor(0xE7, 0x07, 0x11)
RED_DARK = RGBColor(0x8B, 0x1A, 0x1A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x99, 0x99, 0x99)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)

# ── Slide dimensions (16:9) ────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# Use blank layout
blank_layout = prs.slide_layouts[6]


# ── Helper Functions ───────────────────────────────────

def set_bg(slide, color):
    """Set solid background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, italic=False,
             color=WHITE, align=PP_ALIGN.LEFT, font_name="Arial"):
    """Add a text box with single-run formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_multiline(slide, lines, left, top, width, height,
                  font_size=16, color=WHITE, line_spacing=1.5,
                  bold=False, font_name="Arial"):
    """Add text box with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(line_spacing * 4)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = font_name
    return txBox


def add_footer(slide, page_num, total=21):
    """Add footer line + page number (Logo-Platzhalter links)."""
    # Footer line
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(6.9), Inches(12.333), Pt(1)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x33)
    shape.line.fill.background()

    # Linker Slot: Logo-Platzhalter (leer, Luke fügt später das Logo ein)
    # Rechts: Seitenzahl / Total
    add_text(slide, f"{page_num} / {total}",
             Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.4),
             font_size=10, color=GRAY, align=PP_ALIGN.RIGHT)


def add_image_safe(slide, path, left, top, width=None, height=None):
    """Add image if file exists."""
    if os.path.exists(path):
        kwargs = {"image_file": path, "left": left, "top": top}
        if width:
            kwargs["width"] = width
        if height:
            kwargs["height"] = height
        return slide.shapes.add_picture(**kwargs)
    return None


def add_card(slide, left, top, width, height, title, body, bg_color=DARK):
    """Add a rounded card with title and body text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    shape.adjustments[0] = 0.05

    add_text(slide, title, left + Inches(0.3), top + Inches(0.2),
             width - Inches(0.6), Inches(0.5),
             font_size=16, bold=True, color=WHITE)
    add_text(slide, body, left + Inches(0.3), top + Inches(0.65),
             width - Inches(0.6), height - Inches(0.85),
             font_size=12, color=LIGHT_GRAY)


def chapter_slide(slide, number, title, subtitle):
    """Create red chapter divider slide."""
    set_bg(slide, RED_DARK)
    # Big number
    add_text(slide, f"{number:02d}", Inches(0.8), Inches(1.0),
             Inches(4), Inches(1.5),
             font_size=72, bold=True, color=RGBColor(0xFF, 0x99, 0x99),
             font_name="Arial Black")
    # Title
    add_text(slide, title, Inches(0.8), Inches(2.8),
             Inches(10), Inches(1.5),
             font_size=48, bold=True, color=WHITE, font_name="Arial Black")
    # Subtitle
    add_text(slide, subtitle, Inches(0.8), Inches(4.5),
             Inches(10), Inches(1),
             font_size=20, italic=False, color=RGBColor(0xFF, 0xCC, 0xCC))


def add_red_badge(slide, text, left, top, width=Inches(2.5), height=Inches(0.5)):
    """Add a red badge/tag."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RED
    shape.line.fill.background()
    shape.adjustments[0] = 0.3
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Arial"


def add_bullet_list(slide, items, left, top, width, height,
                    icon="✕", icon_color=RED, font_size=14):
    """Add a list with colored icons."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(8)

        icon_run = p.add_run()
        icon_run.text = f"{icon}  "
        icon_run.font.size = Pt(font_size)
        icon_run.font.color.rgb = icon_color
        icon_run.font.bold = True
        icon_run.font.name = "Arial"

        text_run = p.add_run()
        text_run.text = item
        text_run.font.size = Pt(font_size)
        text_run.font.color.rgb = LIGHT_GRAY
        text_run.font.name = "Arial"


def add_check_list(slide, items, left, top, width, height, font_size=14):
    """Add a list with green check icons."""
    add_bullet_list(slide, items, left, top, width, height,
                    icon="✓", icon_color=RGBColor(0x4C, 0xAF, 0x50),
                    font_size=font_size)


# ══════════════════════════════════════════════════════
#  SLIDES
# ══════════════════════════════════════════════════════

# ── SLIDE 1: Titel ─────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, BLACK)
add_text(s, "Fitness Factory Hattingen",
         Inches(0.8), Inches(1.5), Inches(11), Inches(1.5),
         font_size=52, bold=True, font_name="Arial Black")
add_text(s, "Website-Neuauftritt",
         Inches(0.8), Inches(3.2), Inches(11), Inches(1),
         font_size=36, bold=False, color=RED)
add_text(s, "Von der digitalen Visitenkarte\nzum Mitglieder-Magneten.",
         Inches(0.8), Inches(4.5), Inches(11), Inches(1.2),
         font_size=20, color=GRAY)
add_footer(s, 1)


# ── SLIDE 2: Die eine Frage ────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, BLACK)
add_text(s, "Die eine Frage.",
         Inches(0.8), Inches(1.2), Inches(11), Inches(1),
         font_size=48, bold=True, font_name="Arial Black")

# Quote card
shape = s.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.8), Inches(2.8), Inches(11.5), Inches(2.5)
)
shape.fill.solid()
shape.fill.fore_color.rgb = DARK
shape.line.fill.background()
shape.adjustments[0] = 0.03

add_text(s, "Hilft deine Website dabei,\nneue Mitglieder zu gewinnen?",
         Inches(1.3), Inches(3.0), Inches(10.5), Inches(1.5),
         font_size=32, bold=True)
add_text(s, "Wenn nein, arbeitet sie gegen dich.",
         Inches(1.3), Inches(4.3), Inches(10.5), Inches(0.6),
         font_size=18, italic=True, color=RED)

add_text(s, "Jedes Mitglied, das den Kursplan nicht findet,\njede Navigation, die verwirrt — kostet dich Geld.",
         Inches(0.8), Inches(5.8), Inches(11), Inches(0.6),
         font_size=14, color=GRAY)
add_text(s, "Die Website ist kein Design-Projekt. Sie ist eine Verkaufsanlage.",
         Inches(0.8), Inches(6.4), Inches(11), Inches(0.4),
         font_size=13, italic=True, color=RED)
add_footer(s, 2)


# ── SLIDE 3: Website heute ─────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, BLACK)
add_text(s, "Website heute",
         Inches(0.8), Inches(0.5), Inches(6), Inches(0.8),
         font_size=40, bold=True, font_name="Arial Black")
add_text(s, "Was der erste Klick heute kostet",
         Inches(0.8), Inches(1.2), Inches(6), Inches(0.4),
         font_size=14, color=RED)

# Screenshot der alten Website (Desktop-Crop, ohne Cookie-Banner)
desktop_old = os.path.join(SCREENSHOTS_V4, "desktop-old-top.png")
if not os.path.exists(desktop_old):
    desktop_old = os.path.join(SCREENSHOTS_OLD, "01-startseite.png")
add_image_safe(s, desktop_old, Inches(0.8), Inches(1.8), width=Inches(6.5))

add_text(s, "Live-Screenshot fitness-factory-hattingen.de",
         Inches(0.8), Inches(6.25), Inches(6.5), Inches(0.3),
         font_size=10, color=GRAY, italic=True)

# Right side: 3 Badges + Conversion-Impact-Pfeile
add_text(s, "3 Schwachstellen",
         Inches(7.8), Inches(1.8), Inches(5), Inches(0.5),
         font_size=18, bold=True, color=RED)

# Badge 1
add_red_badge(s, "Mobile: Logo unsichtbar, Cookie-Wall vor Hero",
              Inches(7.8), Inches(2.5), width=Inches(4.8), height=Inches(0.55))
add_text(s, "→ Besucher sieht erst Datenschutz-Popup, dann leeren Hero.",
         Inches(7.8), Inches(3.1), Inches(4.8), Inches(0.35),
         font_size=11, color=LIGHT_GRAY)

# Badge 2
add_red_badge(s, "One-Pager getarnt als Seitenstruktur",
              Inches(7.8), Inches(3.55), width=Inches(4.8), height=Inches(0.55))
add_text(s, "→ 7 Menüpunkte, davon 4 nur Sprungmarken (#kursplan, #faq, …).",
         Inches(7.8), Inches(4.15), Inches(4.8), Inches(0.35),
         font_size=11, color=LIGHT_GRAY)

# Badge 3
add_red_badge(s, "Trust-Signale komplett versteckt",
              Inches(7.8), Inches(4.6), width=Inches(4.8), height=Inches(0.55))
add_text(s, "→ 4,9 ★ bei 167 Google-Bewertungen — auf der Startseite null Hinweis.",
         Inches(7.8), Inches(5.2), Inches(4.8), Inches(0.55),
         font_size=11, color=LIGHT_GRAY)

add_text(s, "Browser-Tab der Live-Seite zeigt nur \"Home –\" — kein Studio-Name, kein Ort.",
         Inches(7.8), Inches(6.0), Inches(4.8), Inches(0.5),
         font_size=11, italic=True, color=GRAY)
add_footer(s, 3)


# ── SLIDE 4: Kapitel 01 — Schwachstellen ───────────────
s = prs.slides.add_slide(blank_layout)
chapter_slide(s, 1, "Schwachstellen", "Was Besucher heute erleben —\nund warum sie nicht wiederkommen.")
add_footer(s, 4)


# ── SLIDE 5: Navigation & Header ───────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, BLACK)
add_text(s, "Navigation & Header",
         Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         font_size=36, bold=True, font_name="Arial Black")

# Screenshots oben (echte Live-Site, Desktop 1440×900)
add_image_safe(s, os.path.join(SCREENSHOTS_V4, "desktop-old-top.png"),
               Inches(0.8), Inches(1.3), width=Inches(5.7))
add_image_safe(s, os.path.join(SCREENSHOTS_V4, "desktop-new-top.png"),
               Inches(6.8), Inches(1.3), width=Inches(5.7))

# Vorher-Label
add_text(s, "VORHER (Live)", Inches(0.8), Inches(4.9), Inches(3), Inches(0.4),
         font_size=12, bold=True, color=RED)
add_bullet_list(s, [
    "13 Links ohne Hierarchie",
    "Kein CTA im Header",
    "Cookie-Wall blockiert Hero",
    "Keine Telefon-/WhatsApp-Shortcut"
], Inches(0.8), Inches(5.3), Inches(5.7), Inches(1.3), font_size=12)

# Nachher-Label
add_text(s, "NACHHER (Relaunch)", Inches(6.8), Inches(4.9), Inches(4), Inches(0.4),
         font_size=12, bold=True, color=RGBColor(0x4C, 0xAF, 0x50))
add_check_list(s, [
    "6 klare Hauptpunkte + Logo sichtbar",
    "\"Probetraining buchen\" + \"Jetzt anrufen\"",
    "Hero above-the-fold, keine Popups",
    "Sticky CTA beim Scrollen"
], Inches(6.8), Inches(5.3), Inches(5.7), Inches(1.3), font_size=12)

add_text(s, "→ Besucher findet in 2 Sekunden, wonach er sucht.",
         Inches(0.8), Inches(6.6), Inches(11), Inches(0.3),
         font_size=13, italic=True, color=GRAY)
add_footer(s, 5)


# ── SLIDE 6: Kursplan & Bewertungen ────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, BLACK)
add_text(s, "Kursplan & Bewertungen",
         Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         font_size=36, bold=True, font_name="Arial Black")

# Top-left card: Kursplan vorher
add_card(s, Inches(0.8), Inches(1.5), Inches(5.7), Inches(2.3),
         "Kursplan — vorher",
         "Reine Tabelle ohne Filter.\nKeine Kursbeschreibungen.\nKeine Anmelde-Option.\n16 Kurse — keiner wird emotional verkauft.")

# Top-right card: Kursplan nachher
add_card(s, Inches(6.8), Inches(1.5), Inches(5.7), Inches(2.3),
         "Kursplan — nachher",
         "Klickbares Grid mit Trainer-Foto.\nFilter: Cardio · Kraft · Entspannung.\nJeder Kurs mit Beschreibung + CTA.\n→ Aus Tabelle wird Kurs-Erlebnis.")

# Bottom-left card: Bewertungen vorher
add_card(s, Inches(0.8), Inches(4.2), Inches(5.7), Inches(2.3),
         "4,9 ★ bei 167 Bewertungen — versteckt",
         "Google-Bewertungen nirgendwo\nauf der Startseite sichtbar.\nKein Link zum Google-Profil.\n→ Stärkster Beweis komplett ungenutzt.")

# Bottom-right card: Bewertungen nachher
add_card(s, Inches(6.8), Inches(4.2), Inches(5.7), Inches(2.3),
         "4,9 ★ — jetzt auf jeder Seite",
         "Trust-Bar mit echten Stimmen.\nGoogle-Logo + Sterne-Widget.\nCounter-Animation beim Scrollen.\n→ 167 stille Bewertungen werden Verkäufer.")
add_footer(s, 6)


# ── SLIDE 7: Details die Vertrauen kosten ──────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, BLACK)
add_text(s, "Details, die Vertrauen kosten.",
         Inches(0.8), Inches(0.35), Inches(11), Inches(0.7),
         font_size=32, bold=True, font_name="Arial Black")
add_text(s, "Was Besucher auf der Live-Site sofort sehen",
         Inches(0.8), Inches(1.05), Inches(11), Inches(0.35),
         font_size=13, color=RED)

items = [
    ("Mobile: Logo unsichtbar", "Auf 68 % der Zugriffe fehlt die Marke im Header — kein Wiedererkennungswert."),
    ("Kursplan nicht aktuell", "Angezeigte Kurse stimmen nicht mit Realität überein — sofortiger Vertrauensbruch."),
    ("Anker-Nav wirkt wie Seiten", "7 Menüpunkte, 4 davon springen nur zu Abschnitten — Besucher verliert Orientierung."),
    ("Keine Route-/Anfahrts-Funktion", "Keine eingebettete Karte, kein „Route planen\"-Button — Weg zum Studio fehlt."),
    ("SEPA-Formular öffentlich", "Bankdaten-Seite ungeschützt im Web — DSGVO- und Sicherheitsrisiko."),
    ("GMX-Mailadresse", "Wirkt hobbyhaft neben Premium-Preisen — Domain-Mail ist Standard."),
    ("Kein Team- / Trainer-Bereich", "Kein Gesicht, kein Name — „familiär\" wird behauptet, nicht gezeigt."),
]

row_h = Inches(0.68)
start_y = Inches(1.55)
for i, (title, desc) in enumerate(items):
    y = start_y + row_h * i
    add_text(s, "⚠", Inches(0.8), y, Inches(0.45), Inches(0.4),
             font_size=14, color=RED)
    add_text(s, title, Inches(1.3), y, Inches(3.8), Inches(0.4),
             font_size=13, bold=True)
    add_text(s, desc, Inches(5.2), y, Inches(7.5), Inches(0.5),
             font_size=11, color=LIGHT_GRAY)

add_text(s, "→ Das Studio ist besser als sein Auftritt. Genau das kostet Mitglieder.",
         Inches(0.8), Inches(6.55), Inches(11), Inches(0.4),
         font_size=12, italic=True, color=GRAY)
add_footer(s, 7)


# ── SLIDE 8: Kapitel 02 — Der Neuauftritt ──────────────
s = prs.slides.add_slide(blank_layout)
chapter_slide(s, 2, "Der Neuauftritt",
              "Gebaut für Google.\nGebaut für Handys.\nGebaut für Kunden.")
add_footer(s, 8)


# ── SLIDE 9: Header & Hero ─────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, BLACK)
add_text(s, "Header & Hero",
         Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         font_size=36, bold=True, font_name="Arial Black")
add_text(s, "Die ersten 3 Sekunden verkaufen.\nDer Rest der Seite bestätigt nur noch.",
         Inches(0.8), Inches(1.2), Inches(5), Inches(0.8),
         font_size=16, color=GRAY)

# New hero screenshot
add_image_safe(s, os.path.join(SCREENSHOTS_NEW, "01-hero.png"),
               Inches(6.5), Inches(0.4), width=Inches(6.3))

add_check_list(s, [
    "Starke Headline: \"Dein Gym in Hattingen\"",
    "TrustBar: 4,9 ★ · 167 Bewertungen · 16 Kurse",
    "Primär-CTA: \"Probetraining sichern\" (Rot)",
    "Sekundär: \"Jetzt anrufen\" (Outline)",
    "Sticky CTA-Bar nach 30% Scroll"
], Inches(0.8), Inches(2.3), Inches(5.5), Inches(3.5), font_size=14)

add_text(s, "→ Claim, Beweis, CTA — alles above the fold.",
         Inches(0.8), Inches(6.6), Inches(11), Inches(0.4),
         font_size=13, italic=True, color=GRAY)
add_footer(s, 9)


# ── SLIDE 10: Leistungen & Preise ──────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, BLACK)
add_text(s, "Leistungen & Preise",
         Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         font_size=36, bold=True, font_name="Arial Black")

# Left: Leistungen
add_card(s, Inches(0.8), Inches(1.5), Inches(5.7), Inches(4.8),
         "Leistungen — neu",
         "Aus generischen Icons werden\nechte Räume und echte Versprechen.\n\n"
         "• Jede Leistung mit echtem Foto\n"
         "• Nutzen-Claims statt Feature-Listen\n"
         "• 8 Kernleistungen im 2×4 Grid\n"
         "• CTA: \"Alle Leistungen live erleben\"")

# Right: Preise
add_card(s, Inches(6.8), Inches(1.5), Inches(5.7), Inches(4.8),
         "Preise — neu",
         "3 Pakete. 1 klarer Favorit.\nPreise werden verkauft, nicht gelistet.\n\n"
         "• Basis (55€) · Standard (45€) · Premium (35€)\n"
         "• \"Beliebteste Wahl\" rot markiert\n"
         "• Fremdgeh-Aktion als Trigger\n"
         "• CTA pro Tarif: \"Jetzt Probetraining\"")

add_text(s, "→ \"Alles inklusive\" wird \"Das bekommst du bei uns — und so fühlt es sich an.\"",
         Inches(0.8), Inches(6.6), Inches(11), Inches(0.4),
         font_size=13, italic=True, color=GRAY)
add_footer(s, 10)


# ── SLIDE 11: Kursplan & FAQ ───────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, BLACK)
add_text(s, "Kursplan & FAQ",
         Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         font_size=36, bold=True, font_name="Arial Black")

# Kursplan screenshot
add_image_safe(s, os.path.join(SCREENSHOTS_NEW, "03-kursplan.png"),
               Inches(0.8), Inches(1.5), width=Inches(5.7))

# FAQ info on right
add_text(s, "FAQ — strukturiert",
         Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.5),
         font_size=20, bold=True)
add_check_list(s, [
    "3 Kategorien: Mitgliedschaft · Kurse · Wellness",
    "Ausklappbar, max. 3 Sätze pro Antwort",
    "Kein Juristen-Deutsch",
    "Abschluss-CTA: \"Jetzt WhatsApp schreiben\"",
    "Versprechen: Antwort in unter 2 Stunden"
], Inches(7.0), Inches(2.2), Inches(5.5), Inches(3.5), font_size=13)

add_text(s, "→ Jede beantwortete Frage ist eine gewonnene Mitgliedschaft.",
         Inches(0.8), Inches(6.6), Inches(11), Inches(0.4),
         font_size=13, italic=True, color=GRAY)
add_footer(s, 11)


# ── SLIDE 12: Kontakt & Footer ─────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, BLACK)
add_text(s, "Kontakt & Footer",
         Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         font_size=36, bold=True, font_name="Arial Black")

# Kontakt screenshot
add_image_safe(s, os.path.join(SCREENSHOTS_NEW, "05-kontakt.png"),
               Inches(0.8), Inches(1.5), width=Inches(5.7))

add_check_list(s, [
    "Google-Maps-Embed mit echtem Standort",
    "\"Route planen\"-Button (iOS / Android)",
    "Öffnungszeiten live: \"Jetzt geöffnet · schließt um 22:00\"",
    "WhatsApp-Shortcut direkt daneben",
    "Footer als zweite Startseite mit allen Infos"
], Inches(7.0), Inches(1.8), Inches(5.5), Inches(3.5), font_size=14)

add_text(s, "→ Aus \"Weitere Fragen?\" wird \"Ich komme jetzt vorbei.\"",
         Inches(0.8), Inches(6.6), Inches(11), Inches(0.4),
         font_size=13, italic=True, color=GRAY)
add_footer(s, 12)


# ── SLIDE 13: Mobile-First ─────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, BLACK)
add_text(s, "Mobile-First",
         Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
         font_size=36, bold=True, font_name="Arial Black")
add_text(s, "68 % der Besucher kommen mobil. Ihr erster Eindruck entscheidet alles.",
         Inches(0.8), Inches(1.1), Inches(11), Inches(0.4),
         font_size=15, color=GRAY)

# Mobile Screenshots (echte Live-Shots, iPhone 13 = 1170×2532, width 390)
# Höhe begrenzt auf 4.3" → Breite automatisch proportional
mob_h = Inches(4.3)
# Vorher: mitte-links positionieren
add_image_safe(s, os.path.join(SCREENSHOTS_V4, "mobile-old-top.png"),
               Inches(2.3), Inches(1.7), height=mob_h)
# Nachher: mitte-rechts positionieren
add_image_safe(s, os.path.join(SCREENSHOTS_V4, "mobile-new-top.png"),
               Inches(8.5), Inches(1.7), height=mob_h)

# Labels unter den Screenshots
add_text(s, "VORHER", Inches(1.8), Inches(6.05), Inches(3.5), Inches(0.35),
         font_size=12, bold=True, color=RED, align=PP_ALIGN.CENTER)
add_text(s, "Cookie-Wall, kein Logo, kein CTA",
         Inches(1.8), Inches(6.35), Inches(3.5), Inches(0.3),
         font_size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_text(s, "NACHHER", Inches(8.0), Inches(6.05), Inches(3.5), Inches(0.35),
         font_size=12, bold=True, color=RGBColor(0x4C, 0xAF, 0x50),
         align=PP_ALIGN.CENTER)
add_text(s, "Logo · Hero · Probetraining · Anruf · WhatsApp",
         Inches(8.0), Inches(6.35), Inches(3.5), Inches(0.3),
         font_size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_text(s, "→ Aus \"Cookie wegklicken, dann suchen\" wird \"Sehen, klicken, buchen.\"",
         Inches(0.8), Inches(6.75), Inches(11), Inches(0.3),
         font_size=12, italic=True, color=GRAY, align=PP_ALIGN.CENTER)
add_footer(s, 13)


# ── SLIDE 14: Kapitel 03 — Mehr als eine Website ───────
s = prs.slides.add_slide(blank_layout)
chapter_slide(s, 3, "Mehr als\neine Website",
              "Was zusätzlich im Paket steckt.")
add_footer(s, 14)


# ── SLIDE 15: Dashboard ────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, BLACK)
add_text(s, "Dashboard",
         Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         font_size=36, bold=True, font_name="Arial Black")
add_text(s, "Texte, Kurse, Öffnungszeiten — alles selbst ändern.\nOhne Webentwickler. Ohne Wartezeit.",
         Inches(0.8), Inches(1.2), Inches(11), Inches(0.8),
         font_size=18, color=GRAY)

# Dashboard screenshot from original slides
add_image_safe(s, os.path.join(FF_SLIDES, "FitnessFactory.030.jpeg"),
               Inches(0.8), Inches(2.3), width=Inches(11.5))

add_text(s, "→ Die Website wird zur Plattform. Das Team behält die Kontrolle.",
         Inches(0.8), Inches(6.6), Inches(11), Inches(0.4),
         font_size=13, italic=True, color=GRAY)
add_footer(s, 15)


# ── SLIDE 16: Studio-Auftreten vorher (Klebeband-Zettel) ──
s = prs.slides.add_slide(blank_layout)
set_bg(s, BLACK)
add_text(s, "Nicht nur digital. Auch vor Ort.",
         Inches(0.8), Inches(0.4), Inches(12), Inches(0.8),
         font_size=32, bold=True, font_name="Arial Black")
add_text(s, "Neue Geräte, neue Duschen, neuer Boden — Hinweise aber noch mit Klebeband.",
         Inches(0.8), Inches(1.15), Inches(12), Inches(0.4),
         font_size=15, color=GRAY)

# Drei Vorher-Bilder nebeneinander
img_w = Inches(3.8)
img_h = Inches(4.0)
positions = [Inches(0.55), Inches(4.75), Inches(8.95)]
images_before = [
    (os.path.join(SCHILDER, "Brustpresse", "IMG_2547.HEIC"),
     os.path.join(SCREENSHOTS_V4, "brustpresse_voher.jpg"),
     "Brustpresse",
     "Loser Zettel in einer Folie"),
    (os.path.join(SCHILDER, "Sauna", "Sauna_voher.HEIC"),
     os.path.join(SCREENSHOTS_V4, "sauna_voher.jpg"),
     "Sauna",
     "Papier auf Metallschild geklebt"),
    (os.path.join(SCHILDER, "Spinde", "Spinde_voher.HEIC"),
     os.path.join(SCREENSHOTS_V4, "spinde_voher.jpg"),
     "Spinde",
     "Hinweise handgeklebt"),
]
for x, (_, path, label, hint) in zip(positions, images_before):
    add_image_safe(s, path, x, Inches(1.65), width=img_w, height=img_h)
    add_text(s, label, x, Inches(5.75), img_w, Inches(0.4),
             font_size=16, bold=True, align=PP_ALIGN.CENTER, color=RED)
    add_text(s, hint, x, Inches(6.15), img_w, Inches(0.35),
             font_size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_text(s, "→ Premium-Preis verlangt Premium-Auftreten — auch an der Wand.",
         Inches(0.8), Inches(6.6), Inches(11.5), Inches(0.3),
         font_size=13, italic=True, color=GRAY, align=PP_ALIGN.CENTER)
add_footer(s, 16)


# ── SLIDE 17: Schilder-System im CI-Design (Nachher) ──
s = prs.slides.add_slide(blank_layout)
set_bg(s, BLACK)
add_text(s, "Ein Schilder-System im CI-Design.",
         Inches(0.8), Inches(0.4), Inches(12), Inches(0.8),
         font_size=32, bold=True, font_name="Arial Black")
add_text(s, "Konsistent. Markenkonform. Produktionsreif — für heutige und zukünftige Hinweise.",
         Inches(0.8), Inches(1.15), Inches(12), Inches(0.4),
         font_size=15, color=GRAY)

images_after = [
    (os.path.join(SCHILDER, "Brustpresse", "Brustpresse_nachher.jpeg"),
     "Brustpresse", "Geräte-Hinweis im Rahmen"),
    (os.path.join(SCHILDER, "Sauna", "Sauna_nachher.jpeg"),
     "Sauna", "Regeln sauber im Rahmen"),
    (os.path.join(SCHILDER, "Spinde", "Spinde_nacher.jpeg"),
     "Spinde", "CI-Schild mit Logo"),
]
for x, (path, label, hint) in zip(positions, images_after):
    add_image_safe(s, path, x, Inches(1.65), width=img_w, height=img_h)
    add_text(s, label, x, Inches(5.75), img_w, Inches(0.4),
             font_size=16, bold=True, align=PP_ALIGN.CENTER,
             color=RGBColor(0x4C, 0xAF, 0x50))
    add_text(s, hint, x, Inches(6.15), img_w, Inches(0.35),
             font_size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_text(s, "→ Ein Logo, eine Website, ein Studio — alles aus einem Guss.",
         Inches(0.8), Inches(6.6), Inches(11.5), Inches(0.3),
         font_size=13, italic=True, color=GRAY, align=PP_ALIGN.CENTER)
add_footer(s, 17)


# ── SLIDE 18: Der Preis des Nichtstuns ─────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, BLACK)
add_text(s, "Der Preis des Nichtstuns",
         Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
         font_size=36, bold=True, font_name="Arial Black")
add_text(s, "Was eine halbfertige Website jeden Monat kostet.",
         Inches(0.8), Inches(1.25), Inches(11), Inches(0.4),
         font_size=14, color=GRAY)

# 4 cards (3 Zahl-Karten + 1 Mustererkennung)
cards = [
    ("68 %", "kommen mobil", "Und stolpern über\nCookie-Wall + fehlendes Logo."),
    ("3 Sek.", "Entscheidungszeit", "Dann ist der\nBesucher weg."),
    ("~2.430 €", "entgangener Umsatz/Monat", "~90 verlorene Kontakte,\n~5 % davon = Mitglieder."),
    ("3 × / 4 Jahre", "Websites gekauft", "Ohne Verkaufs-Wirkung.\nDas Muster muss enden."),
]
card_w = Inches(2.85)
gap = Inches(0.15)
total = card_w * 4 + gap * 3
start_x = (SLIDE_W - total) / 2
for i, (number, label, detail) in enumerate(cards):
    x = start_x + (card_w + gap) * i
    card = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x, Inches(1.9), card_w, Inches(4.2)
    )
    card.fill.solid()
    # Die letzte Karte (Muster) dunkelrot hervorheben
    card.fill.fore_color.rgb = RED_DARK if i == 3 else DARK
    card.line.fill.background()
    card.adjustments[0] = 0.05

    # Zahl: für lange Zahlen font_size kleiner
    n_size = 36 if len(number) > 6 else 44
    add_text(s, number, x + Inches(0.2), Inches(2.2),
             card_w - Inches(0.4), Inches(1.2),
             font_size=n_size, bold=True,
             color=RGBColor(0xFF, 0xCC, 0xCC) if i == 3 else RED,
             font_name="Arial Black")
    add_text(s, label, x + Inches(0.2), Inches(3.65),
             card_w - Inches(0.4), Inches(0.5),
             font_size=13, bold=True)
    add_text(s, detail, x + Inches(0.2), Inches(4.4),
             card_w - Inches(0.4), Inches(1.4),
             font_size=11, color=LIGHT_GRAY if i == 3 else GRAY)

add_text(s, "Jede Woche ohne verkaufende Website kostet Mitgliedschaften.",
         Inches(0.8), Inches(6.3), Inches(11), Inches(0.4),
         font_size=15, italic=True, color=RED, align=PP_ALIGN.CENTER)
add_footer(s, 18)


# ── SLIDE 19: Farbwelt & Marke ─────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, BLACK)
add_text(s, "Farbwelt & Marke",
         Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         font_size=36, bold=True, font_name="Arial Black")
add_text(s, "Sechs Farben. Ein klares Versprechen.\nDirekt aus dem Logo abgeleitet.",
         Inches(0.8), Inches(1.2), Inches(11), Inches(0.7),
         font_size=16, color=GRAY)

# Color swatches
colors = [
    ("#FFFFFF", "Weiß", "Hintergrund"),
    ("#E70711", "Hellrot", "Buttons, CTA"),
    ("#C70E16", "Rot", "Logo, Akzent"),
    ("#721F01", "Dunkelrot", "Tiefe, Text"),
    ("#300F06", "Braun", "Logo-Detail"),
    ("#000000", "Schwarz", "Headlines"),
]
for i, (hex_c, name, use) in enumerate(colors):
    x = Inches(0.8) + Inches(2.05) * i
    # Color swatch
    swatch = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x, Inches(2.2), Inches(1.8), Inches(1.8)
    )
    r = int(hex_c[1:3], 16)
    g = int(hex_c[3:5], 16)
    b = int(hex_c[5:7], 16)
    swatch.fill.solid()
    swatch.fill.fore_color.rgb = RGBColor(r, g, b)
    if hex_c == "#FFFFFF":
        swatch.line.color.rgb = GRAY
    else:
        swatch.line.fill.background()
    swatch.adjustments[0] = 0.08

    add_text(s, name, x, Inches(4.2), Inches(1.8), Inches(0.4),
             font_size=13, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, use, x, Inches(4.55), Inches(1.8), Inches(0.4),
             font_size=11, color=GRAY, align=PP_ALIGN.CENTER)

# Logo note
add_text(s, "Logo-Empfehlung:",
         Inches(0.8), Inches(5.3), Inches(11), Inches(0.4),
         font_size=14, bold=True, color=RED)
add_text(s, "Blitz vereinfachen oder durch neues Symbol ersetzen — \"ITNESS\" statt \"FITNESS\" liest sich falsch.\nFlach, modern, skalierbar — wie FitX, John Reed, Urban Sports Club.",
         Inches(0.8), Inches(5.7), Inches(11), Inches(0.8),
         font_size=13, color=LIGHT_GRAY)
add_footer(s, 19)


# ── SLIDE 20: Nächste Schritte ─────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, BLACK)
add_text(s, "Nächste Schritte",
         Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         font_size=40, bold=True, font_name="Arial Black")
add_text(s, "Vom Ist-Zustand zum Neuauftritt — in 4 Wochen.",
         Inches(0.8), Inches(1.2), Inches(7), Inches(0.5),
         font_size=16, color=GRAY)

steps = [
    ("1", "Freigabe Farbwelt & Typografie", "Woche 1"),
    ("2", "Logo-Neugestaltung & Bildbearbeitung", "Woche 1–2"),
    ("3", "Website-Relaunch mit neuem Kursplan", "Woche 2–3"),
    ("4", "Instagram-Update & Highlight-Covers", "Woche 3–4"),
    ("5", "Live-Gang & Feinschliff", "Woche 4"),
]

for i, (num, title, week) in enumerate(steps):
    y = Inches(2.0) + Inches(0.9) * i
    # Red circle with number
    circle = s.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(0.8), y, Inches(0.55), Inches(0.55)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = RED
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = num
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Arial"

    add_text(s, title, Inches(1.6), y + Inches(0.02),
             Inches(6), Inches(0.35),
             font_size=17, bold=True)
    add_text(s, week, Inches(1.6), y + Inches(0.35),
             Inches(6), Inches(0.3),
             font_size=12, color=GRAY)

# "Bereit?" card on right
bereit = s.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(8.5), Inches(2.5), Inches(4), Inches(3.5)
)
bereit.fill.solid()
bereit.fill.fore_color.rgb = RED_DARK
bereit.line.fill.background()
bereit.adjustments[0] = 0.05

add_text(s, "Bereit?", Inches(9.0), Inches(2.8),
         Inches(3), Inches(0.8),
         font_size=36, bold=True, font_name="Arial Black")
add_text(s, "Nach Freigabe:\n7 Tage bis Live-Gang.",
         Inches(9.0), Inches(3.7), Inches(3), Inches(1),
         font_size=16, color=RGBColor(0xFF, 0xCC, 0xCC))
add_text(s, "Lass uns starten.",
         Inches(9.0), Inches(4.9), Inches(3), Inches(0.4),
         font_size=18, bold=True, italic=True, color=WHITE)
add_text(s, "Diesmal gebaut,\num zu bleiben.",
         Inches(9.0), Inches(5.4), Inches(3), Inches(0.8),
         font_size=13, italic=True, color=RGBColor(0xFF, 0xCC, 0xCC))

add_footer(s, 20)


# ── SLIDE 21: Abschluss ───────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, BLACK)
add_text(s, "Die Website ist fertig.",
         Inches(0.8), Inches(2.0), Inches(11), Inches(1.2),
         font_size=52, bold=True, font_name="Arial Black")
add_text(s, "Sie wartet auf dein Go.",
         Inches(0.8), Inches(3.5), Inches(11), Inches(1),
         font_size=28, color=RED)
add_text(s, "Luke Kozik\nWeb Development & Digital Strategy",
         Inches(0.8), Inches(5.5), Inches(5), Inches(1),
         font_size=14, color=GRAY)
add_footer(s, 21)


# ── Save ───────────────────────────────────────────────
prs.save(OUTPUT)
print(f"✓ Präsentation erstellt: {OUTPUT}")
print(f"  → {len(prs.slides)} Slides")
